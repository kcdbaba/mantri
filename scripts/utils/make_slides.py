#!/usr/bin/env python3
"""
make_slides.py

Converts staff_interview_presentation.md into a .pptx file using Gemini
to extract clean bullet points from the ASCII art slide content.

Upload the output .pptx to Google Drive — it will auto-convert to Google Slides.

Usage:
    python scripts/make_slides.py
    python scripts/make_slides.py --input interviews/staff_interview_presentation.md
    python scripts/make_slides.py --output demo/staff_interview.pptx

Setup:
    pip install python-pptx google-genai python-dotenv
    GOOGLE_API_KEY must be set in .env
"""

import argparse
import os
import re
import sys
from pathlib import Path

from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

load_dotenv()

# ── Colours ──────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy background
ACCENT    = RGBColor(0x4A, 0x90, 0xD9)   # blue accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Parse markdown into slide blocks ─────────────────────────────────────────

def parse_slides(md_text: str) -> list[dict]:
    """
    Returns list of dicts with keys:
        number   int
        title    str
        raw      str   (ASCII art block only — for Gemini content extraction)
        script   str   (speaker script from **SCRIPT:** section)
    """
    pattern = re.compile(r'^## (SLIDE\s+\d+\s*[—–-]+\s*.+)$', re.MULTILINE)
    positions = [(m.start(), m.group(1)) for m in pattern.finditer(md_text)]

    slides = []
    for i, (pos, heading) in enumerate(positions):
        end = positions[i + 1][0] if i + 1 < len(positions) else len(md_text)
        section = md_text[pos:end].strip()

        num_match = re.search(r'SLIDE\s+(\d+)', heading)
        number = int(num_match.group(1)) if num_match else i + 1

        title_match = re.search(r'SLIDE\s+\d+\s*[—–-]+\s*(.+)', heading)
        title = title_match.group(1).strip() if title_match else heading

        # Split off the script section
        script_match = re.search(r'\*\*SCRIPT:\*\*\s*\n(.*)', section, re.DOTALL)
        script = script_match.group(1).strip() if script_match else ""

        # Raw = everything before **SCRIPT:** (the ASCII art block)
        raw = section[:script_match.start()].strip() if script_match else section

        slides.append({"number": number, "title": title, "raw": raw, "script": script})

    return slides


# ── Gemini: extract structured content ───────────────────────────────────────

EXTRACT_PROMPT = """\
You are formatting a staff presentation about an AI operations agent called Mantri.

Below is the raw content of one slide (ASCII art box layout only).

Extract the slide content as clean structured text. Return ONLY a JSON object with these fields:
- "title": short slide title (5 words max, no "SLIDE N" prefix)
- "bullets": list of 3-6 bullet point strings (concise, plain text, no markdown symbols)

Rules:
- Extract content from inside the ASCII art box
- Keep bullets short and readable on a slide (under 12 words each)
- Do not include facilitator instructions like "[Leave this on screen]"
- For interview question slides, the first bullet should be the question itself, followed by any sub-points
- Preserve the meaning and intent of every line — do not merge or drop content

Raw slide content:
{raw}
"""

def extract_slide_content(client: genai.Client, raw: str) -> dict:
    import json
    prompt = EXTRACT_PROMPT.format(raw=raw)
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=types.Content(role="user", parts=[types.Part(text=prompt)]),
        config=types.GenerateContentConfig(
            response_mime_type="application/json",
        ),
    )
    result = json.loads(response.text)
    # Ensure "bullets" key exists
    result.setdefault("bullets", [])
    return result


# ── Build PPTX ────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_DARK)

    # Title
    txb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "MANTRI — Operations Assistant"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    txb2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(0.8))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "A conversation about how we work together"
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT

    # Footer
    txb3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10), Inches(0.6))
    tf3 = txb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Kunal Chowdhury  ·  Ashish Chhabra  ·  Guwahati"
    p3.font.size = Pt(16)
    p3.font.color.rgb = LIGHT


def add_content_slide(prs: Presentation, number: int, title: str, bullets: list[str], script: str = ""):
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_DARK)

    # Slide number badge
    txb_num = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(0.8), Inches(0.5))
    tf_num = txb_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = str(number)
    p_num.font.size = Pt(13)
    p_num.font.color.rgb = ACCENT

    # Title
    txb_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.25), Inches(11.5), Inches(0.9))
    tf_title = txb_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(30)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE

    # Accent line under title
    line = slide.shapes.add_shape(
        1,
        Inches(1.0), Inches(1.2), Inches(11.3), Emu(35000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Bullet body
    txb_body = slide.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.3), Inches(5.7))
    tf_body = txb_body.text_frame
    tf_body.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        run = p.add_run()
        run.text = f"\u2022  {bullet}"   # • bullet character
        run.font.size = Pt(21)
        run.font.color.rgb = LIGHT
        p.space_before = Pt(10)
        p.space_after = Pt(4)

    # Speaker notes — full script from markdown
    if script:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = script


def build_pptx(slides_data: list[dict], out_path: Path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 is title — use custom layout
    add_title_slide(prs)

    # Remaining slides
    for s in slides_data[1:]:
        add_content_slide(
            prs,
            number=s["number"],
            title=s["content"]["title"],
            bullets=s["content"]["bullets"],
            script=s["script"],
        )

    prs.save(str(out_path))
    print(f"Saved: {out_path}")


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--input",  default="interviews/staff_interview_presentation.md")
    parser.add_argument("--output", default="demo/staff_interview.pptx")
    args = parser.parse_args()

    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        print("GOOGLE_API_KEY not set"); sys.exit(1)

    client = genai.Client(api_key=api_key)

    md_path = Path(args.input)
    if not md_path.exists():
        print(f"Not found: {md_path}"); sys.exit(1)

    md_text = md_path.read_text(encoding="utf-8")
    raw_slides = parse_slides(md_text)
    print(f"Found {len(raw_slides)} slides")

    slides_data = []
    for s in raw_slides:
        print(f"  Extracting slide {s['number']}: {s['title'][:50]}")
        content = extract_slide_content(client, s["raw"])
        slides_data.append({**s, "content": content})

    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    build_pptx(slides_data, out_path)

    print(f"\nDone. Upload {out_path} to Google Drive to convert to Google Slides.")


if __name__ == "__main__":
    main()
